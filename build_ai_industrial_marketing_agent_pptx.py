# -*- coding: utf-8 -*-
# 依赖：pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ===== 可自定义品牌区 =====
TITLE = "AI工业营销Agent — 产品定义与方案说明（10页）"
SUBTITLE = "从“想法/素材”到“成片/成稿”的工业级营销内容自动化引擎"
DATE = "2026-01-13"
VERSION = "v1.0"
AUTHOR = "yy1 / 深思"
LOGO_PATH = "assets/logo.png"  # 请将公司LOGO保存为 assets/logo.png
BRAND_FONT = "Microsoft YaHei"
PRIMARY = RGBColor(16, 54, 107)   # 主色深蓝
ACCENT = RGBColor(0, 160, 233)    # 强调色蓝
TEXT_COLOR = RGBColor(30, 30, 30)

# ===== 工具函数 =====
def set_text_style(run, bold=False, size=20, color=TEXT_COLOR, font=BRAND_FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_page_number(slide, idx, total):
    left = Inches(9.0)
    top = Inches(6.9)
    width = Inches(1.0)
    height = Inches(0.3)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{idx}/{total}"
    set_text_style(run, size=10, color=RGBColor(120, 120, 120))
    p.alignment = PP_ALIGN.RIGHT

def add_logo(slide):
    try:
        slide.shapes.add_picture(LOGO_PATH, Inches(9.0), Inches(0.2), height=Inches(0.5))
    except Exception:
        # 当未提供 LOGO 时，忽略
        pass

def add_title(slide, title):
    title_shape = slide.shapes.title
    if title_shape is None:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.0), Inches(1.0))
        tf = txBox.text_frame
    else:
        tf = title_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_text_style(r, bold=True, size=28, color=PRIMARY)

def add_bullets(slide, items, left=Inches(0.8), top=Inches(1.8), width=Inches(8.8), height=Inches(4.8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()

    # 为了首段也能正常显示，先创建一个空段落
    first_p = tf.paragraphs[0]
    first_p.text = ""

    def add_item(text, level=0):
        p = tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        set_text_style(run, size=18 if level == 0 else 16, color=TEXT_COLOR)

    for item in items:
        if isinstance(item, tuple):
            add_item(item[0], item[1])
        elif isinstance(item, dict):
            add_item(item.get("text", ""), item.get("level", 0))
        else:
            add_item(str(item), 0)

def add_cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    # 背景条形装饰
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.5), Inches(10), Inches(0.4))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    # 标题
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.8), Inches(1.2))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = TITLE
    set_text_style(r, bold=True, size=36, color=PRIMARY)
    # 副标题
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = SUBTITLE
    set_text_style(r, size=20, color=TEXT_COLOR)
    # 版本与作者
    tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(6.5), Inches(1.0))
    tf2 = tx2.text_frame
    tf2.clear()
    p = tf2.paragraphs[0]
    r = p.add_run()
    r.text = f"日期：{DATE}    版本：{VERSION}    作者：{AUTHOR}"
    set_text_style(r, size=14, color=RGBColor(90, 90, 90))
    add_logo(slide)
    add_page_number(slide, 1, 10)

def add_simple_slide(prs, idx, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    add_title(slide, title)
    add_bullets(slide, bullets)
    add_logo(slide)
    add_page_number(slide, idx, 10)

def build_deck():
    prs = Presentation()
    # 1 封面
    add_cover_slide(prs)

    # 2 愿景与定位 + 核心价值
    add_simple_slide(prs, 2, "愿景与定位 + 核心价值", [
        "愿景：成为工业企业“营销内容生产力中枢”",
        "定位：面向工业的多Agent编排平台，覆盖视频/海报/宣传册",
        "核心价值：",
        ("速度：从“想法+素材”到“可投放物料”分钟级", 1),
        ("专业：工业知识图谱 + LLM 推理，合规表达", 1),
        ("一致性：品牌规范/术语/尺寸与配色统一", 1),
        ("降本增效：减少外包与反复沟通成本", 1),
        ("优化闭环：A/B 测试与投放数据驱动", 1),
    ])

    # 3 背景与痛点
    add_simple_slide(prs, 3, "背景与痛点", [
        "行业复杂：参数/标准/工况/法规差异大，内容易“外行化”",
        "周期冗长：方案撰写、物料制作、审校耗时",
        "工具分散：多供应商工具，资产难沉淀",
        "闭环弱：投放效果难反哺内容优化",
        "合规风险：用词/承诺/版权/认证引用易出错",
    ])

    # 4 目标用户与应用场景
    add_simple_slide(prs, 4, "目标用户与应用场景", [
        "目标用户：市场/品牌/公关、售前工程、渠道团队、代理商",
        "典型场景：",
        ("新品发布：30s亮点视频 + KV海报 + 4P宣传册", 1),
        ("方案营销：工况/行业应用视频与单页", 1),
        ("展会速产：展前一周批量物料生成", 1),
        ("渠道共创：伙伴素材输入，统一风格输出", 1),
        ("海外本地化：多语言字幕/配音/度量单位/法规适配", 1),
    ])

    # 5 核心能力与输出类型
    add_simple_slide(prs, 5, "核心能力与输出类型", [
        "输入：想法/brief、参数表、素材库（CAD/渲染/视频）、案例文档",
        "输出：广告视频（15/30/60s）、多规格海报、宣传册/白皮书",
        "智能增强：",
        ("要点提炼、痛点-卖点映射、脚本/镜头联动", 1),
        ("配音/字幕/动效建议与套版", 1),
        ("品牌与合规校验；A/B版本与效果预测", 1),
    ])

    # 6 端到端体验与Demo要点
    add_simple_slide(prs, 6, "端到端体验与Demo要点", [
        "流程：选模板 → 输入想法/素材 → 初稿预览 → 多版本生成 → 校验 → 导出/投放",
        "视频故事板（30s）：",
        ("场景（痛点）→ 卖点（性能/可靠性/TCO） → 证据（数据/认证） → CTA", 1),
        "海报/宣传册：标题/主视觉/卖点与参数/证据区/CTA/品牌规范自动套用",
    ])

    # 7 技术架构与多模型生态 + 协议
    add_simple_slide(prs, 7, "技术架构与多模型生态 + 协议", [
        "Agent中枢：任务编排、意图识别、技能路由、版本管理",
        "语��与推理：自研工业LLM、gemini3-pro、GPT5（规划）、deepseek",
        "视觉/视频：可灵、火山、Veo3、sora2、通义万相、Vidu、海螺",
        "协议：Function Call、MCP（Model Context Protocol）、A2A（Agent-to-Agent）",
        "适配层：能力声明、成本/时延画像、质量档案、故障转移/降级",
    ])

    # 8 工业知识图谱与数据治理 + 合规
    add_simple_slide(prs, 8, "工业知识图谱与数据治理 + 合规", [
        "知识图谱：行业→工艺→设备→部件→材料→参数→标准→认证→案例",
        "用途：卖点生成、术语统一、选型逻辑、合规守则",
        "数据治理：资产/模板库、权限与水印、检索与生命周期",
        "合规：禁用词与承诺用语、认证引用格式、版权许可与溯源",
    ])

    # 9 商业模式、KPI与路线图
    add_simple_slide(prs, 9, "商业模式、KPI与路线图", [
        "商业模式：订阅；企业版（私有化/SLA/定制）；生态分成；增值服务（审校/品牌/数据治理）",
        "KPI：TTV↓≥70%；外包成本↓≥40%；一致性≥90；合规≥98%；投放指标提升",
        "路线图：",
        ("Q1：核心闭环与2家视觉工具打通", 1),
        ("Q2：多Agent稳定、A2A跨系统、合规1.0", 1),
        ("Q3：多语言本地化、自动AB测试、≥5工具接入", 1),
        ("Q4：行业包、私有化SLA、生态伙伴计划", 1),
    ])

    # 10 下一步与CTA
    add_simple_slide(prs, 10, "下一步与CTA", [
        "预约“从Brief到导出”的实机演示",
        "选定试点产品线与素材，一周交付首批物料",
        "提供品牌手册/LOGO/色板，生成品牌化模板",
        "如需我直接上传PPTX至GitHub，请提供 owner/repo 与目标分支",
    ])

    return prs

if __name__ == "__main__":
    deck = build_deck()
    output_name = "AI工业营销Agent-产品定义_10页版.pptx"
    deck.save(output_name)
    print(f"已生成：{output_name}")